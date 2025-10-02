#!/usr/bin/env python3
"""
Quick script to check PowerPoint content with word counts
"""
from pptx import Presentation

prs = Presentation('output/Effect_of_Dexamethasone_on_Ventilator-Free_Days_in.pptx')

# Field limits
limits = {
    'Population': 15,
    'Intervention': 15,
    'Setting': 10,
    'Primary Outcome': 20,
    'Finding 1': 15,
    'Finding 2': 15
}

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_num}")
    print('='*60)

    current_field = None
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()

            # Check if this is a field label
            if text in limits:
                current_field = text
                print(f"\n[{text}]")
            elif current_field and text:
                # This is the content for the current field
                word_count = len(text.split())
                limit = limits.get(current_field, 999)
                status = "✅" if word_count <= limit else f"❌ OVER LIMIT!"
                print(f"  {status} ({word_count}/{limit} words): {text[:100]}...")
                current_field = None
            else:
                # Other content (title, footer, etc)
                if len(text) > 20:  # Skip short labels
                    print(f"\n{text}")
