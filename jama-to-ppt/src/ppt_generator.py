"""
VA Format PowerPoint Generator
"""

from typing import Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class VAPowerPointGenerator:
    """
    Generates PowerPoint presentations in Veterans Affairs (VA) format
    """

    # VA Color Palette
    COLORS = {
        'background': RGBColor(0, 51, 102),      # #003366 - Dark blue
        'title_text': RGBColor(255, 255, 255),   # #FFFFFF - White
        'box_background': RGBColor(255, 255, 255),  # #FFFFFF - White
        'box_border': RGBColor(0, 51, 102),      # #003366 - VA blue
        'box_title': RGBColor(0, 51, 102),       # #003366 - VA blue
        'box_content': RGBColor(0, 0, 0),        # #000000 - Black
        'footer_text': RGBColor(200, 200, 200),  # #C8C8C8 - Light gray
    }

    # Dimensions (in inches)
    SLIDE_WIDTH = 10
    SLIDE_HEIGHT = 7.5

    # Title box
    TITLE_WIDTH = 9
    TITLE_HEIGHT = 1
    TITLE_LEFT = 0.5
    TITLE_TOP = 0.3

    # Info boxes
    BOX_WIDTH = 3.5
    BOX_HEIGHT = 1.2

    # Columns
    LEFT_COLUMN_X = 1.5
    RIGHT_COLUMN_X = 5.5

    # Rows
    ROW_1_Y = 1.8
    ROW_2_Y = 3.5
    ROW_3_Y = 5.2

    # Footer
    FOOTER_Y = 7
    FOOTER_WIDTH = 9

    def __init__(self, article_data: Dict[str, str], icon_type: str = 'medical', verbose: bool = False):
        self.data = article_data
        self.icon_type = icon_type
        self.verbose = verbose
        self.prs = Presentation()

        # Set slide dimensions
        self.prs.slide_width = Inches(self.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT)

    def generate(self, output_path: str):
        """Generate PowerPoint presentation"""
        if self.verbose:
            print("📝 PowerPoint oluşturuluyor...")

        # Add blank slide
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']

        # Add title
        self._add_title(slide)

        # Add info boxes
        self._add_info_box(slide, "Population", self.data['population'],
                          self.LEFT_COLUMN_X, self.ROW_1_Y)
        self._add_info_box(slide, "Intervention", self.data['intervention'],
                          self.RIGHT_COLUMN_X, self.ROW_1_Y)
        self._add_info_box(slide, "Setting", self.data['setting'],
                          self.LEFT_COLUMN_X, self.ROW_2_Y)
        self._add_info_box(slide, "Primary Outcome", self.data['primary_outcome'],
                          self.RIGHT_COLUMN_X, self.ROW_2_Y)
        self._add_info_box(slide, "Finding 1", self.data['finding_1'],
                          self.LEFT_COLUMN_X, self.ROW_3_Y)
        self._add_info_box(slide, "Finding 2", self.data['finding_2'],
                          self.RIGHT_COLUMN_X, self.ROW_3_Y)

        # Add footer
        self._add_footer(slide)

        # Save presentation
        self.prs.save(output_path)

        if self.verbose:
            print(f"✅ PowerPoint kaydedildi: {output_path}")

    def _add_title(self, slide):
        """Add title box to slide"""
        title = self.data['title']

        # Determine font size based on title length
        if len(title) > 80:
            font_size = 18
        else:
            font_size = 24

        # Create title text box
        left = Inches(self.TITLE_LEFT)
        top = Inches(self.TITLE_TOP)
        width = Inches(self.TITLE_WIDTH)
        height = Inches(self.TITLE_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add text
        p = text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER

        # Format text
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = self.COLORS['title_text']

    def _add_info_box(self, slide, box_title: str, content: str, left_pos: float, top_pos: float):
        """Add information box with title and content"""
        left = Inches(left_pos)
        top = Inches(top_pos)
        width = Inches(self.BOX_WIDTH)
        height = Inches(self.BOX_HEIGHT)

        # Create box shape
        box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )

        # Format box
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['box_background']
        box.line.color.rgb = self.COLORS['box_border']
        box.line.width = Pt(2)

        # Add text frame
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)

        # Clear default paragraph
        text_frame.clear()

        # Add title paragraph
        title_p = text_frame.paragraphs[0]
        title_p.text = box_title
        title_p.alignment = PP_ALIGN.LEFT
        title_run = title_p.runs[0]
        title_run.font.bold = True
        title_run.font.size = Pt(12)
        title_run.font.color.rgb = self.COLORS['box_title']

        # Add content paragraph
        content_p = text_frame.add_paragraph()
        content_p.text = content
        content_p.alignment = PP_ALIGN.LEFT

        # Try different font sizes if text overflows
        font_size = 10
        content_run = content_p.runs[0]
        content_run.font.size = Pt(font_size)
        content_run.font.color.rgb = self.COLORS['box_content']

        # Check if text fits, reduce font size if needed
        if self._text_overflows(content, width, height):
            for size in [9, 8]:
                content_run.font.size = Pt(size)
                if not self._text_overflows(content, width, height):
                    break
            else:
                # If still overflows, truncate with ellipsis
                max_chars = int(len(content) * 0.7)
                content_p.text = content[:max_chars] + '...'

    def _text_overflows(self, text: str, width, height) -> bool:
        """
        Simple heuristic to check if text might overflow
        This is approximate - actual overflow depends on font metrics
        """
        # Rough estimate: ~10 characters per inch at 10pt
        chars_per_line = int(width.inches * 10)
        approx_lines = len(text) / chars_per_line
        max_lines = int(height.inches * 5)  # ~5 lines per inch

        return approx_lines > max_lines

    def _add_footer(self, slide):
        """Add footer with authors, date, and DOI"""
        footer_text = f"{self.data['authors']} | {self.data['publication_date']} | DOI: {self.data['doi']}"

        left = Inches(0.5)
        top = Inches(self.FOOTER_Y)
        width = Inches(self.FOOTER_WIDTH)
        height = Inches(0.4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.CENTER

        run = p.runs[0]
        run.font.size = Pt(9)
        run.font.color.rgb = self.COLORS['footer_text']
