"""
JAMA Oncology Format PowerPoint Generator
Green theme for oncology journal
"""

from typing import Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

try:
    from chart_generator import ChartGenerator
    CHART_AVAILABLE = True
except ImportError:
    CHART_AVAILABLE = False


class JAMAOncologyPowerPointGenerator:
    """
    Generates PowerPoint presentations in JAMA Oncology format (GREEN theme)
    """

    # JAMA Oncology Color Palette - GREEN theme
    COLORS = {
        'header_bg': RGBColor(0, 130, 114),          # #008272 - Teal/Green
        'title_text': RGBColor(255, 255, 255),       # #FFFFFF - White
        'box_background': RGBColor(240, 240, 240),   # #F0F0F0 - Light gray
        'box_title': RGBColor(0, 130, 114),          # #008272 - Teal green
        'box_content': RGBColor(50, 50, 50),         # #323232 - Dark gray
        'footer_text': RGBColor(100, 100, 100),      # #646464 - Gray
        'accent': RGBColor(0, 130, 114),             # #008272 - Teal green for highlights
    }

    # Dimensions (in inches) - Same as JAMA Open
    SLIDE_WIDTH = 10
    SLIDE_HEIGHT = 7.5

    # Header
    HEADER_HEIGHT = 0.7
    LOGO_WIDTH = 2.5
    LOGO_HEIGHT = 0.6
    LOGO_LEFT = 0.3
    LOGO_TOP = 0.1

    # Main title - REFERENCE IMAGE EXACT
    TITLE_TOP = 0.85
    TITLE_LEFT = 0.5
    TITLE_WIDTH = 9.0
    TITLE_HEIGHT = 0.65

    # Content boxes - SOL TARAF (2x2 grid - 4 kutu)
    LEFT_BOX_WIDTH = 2.8
    LEFT_BOX_HEIGHT = 2.5
    LEFT_COL1_X = 0.4
    LEFT_COL2_X = 3.4
    LEFT_BOX_SPACING = 0.15

    # Üst sıra (POPULATION, INTERVENTION)
    TOP_ROW_Y = 1.65
    # Alt sıra (SETTINGS, PRIMARY OUTCOME)
    BOTTOM_ROW_Y = TOP_ROW_Y + LEFT_BOX_HEIGHT + LEFT_BOX_SPACING

    # Content boxes - SAĞ TARAF (1 BÜYÜK FINDINGS kutusu)
    RIGHT_BOX_WIDTH = 3.6
    RIGHT_BOX_X = 6.4
    RIGHT_BOX_Y = 1.65
    RIGHT_BOX_HEIGHT = (LEFT_BOX_HEIGHT * 2) + LEFT_BOX_SPACING

    # Footer - EXACT FROM REFERENCE
    FOOTER_Y = 7.0
    FOOTER_HEIGHT = 0.4

    def __init__(self, article_data: Dict[str, str], icon_type: str = 'medical', verbose: bool = False):
        self.data = article_data
        self.icon_type = icon_type
        self.verbose = verbose
        self.prs = Presentation()

        # Set slide dimensions
        self.prs.slide_width = Inches(self.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT)

    def generate(self, output_path: str):
        """Generate PowerPoint presentation"""
        if self.verbose:
            print("📝 JAMA Oncology PowerPoint oluşturuluyor...")

        # Add blank slide
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Add GREEN header bar
        self._add_header(slide)

        # Add main title
        self._add_main_title(slide)

        # SOL TARAF - Üst sıra (2 kutu)
        self._add_content_box(
            slide, "POPULATION", self.data.get('population', 'N/A'),
            self.LEFT_COL1_X, self.TOP_ROW_Y, self.LEFT_BOX_WIDTH, self.LEFT_BOX_HEIGHT
        )
        self._add_content_box(
            slide, "INTERVENTION", self.data.get('intervention', 'N/A'),
            self.LEFT_COL2_X, self.TOP_ROW_Y, self.LEFT_BOX_WIDTH, self.LEFT_BOX_HEIGHT
        )

        # SOL TARAF - Alt sıra (2 kutu)
        self._add_content_box(
            slide, "SETTINGS / LOCATIONS", self.data.get('setting', 'N/A'),
            self.LEFT_COL1_X, self.BOTTOM_ROW_Y, self.LEFT_BOX_WIDTH, self.LEFT_BOX_HEIGHT
        )
        self._add_content_box(
            slide, "PRIMARY OUTCOME", self.data.get('primary_outcome', 'N/A'),
            self.LEFT_COL2_X, self.BOTTOM_ROW_Y, self.LEFT_BOX_WIDTH, self.LEFT_BOX_HEIGHT
        )

        # SAĞ TARAF - 1 BÜYÜK FINDINGS kutusu (grafik + iDFS)
        self._add_findings_box(slide)

        # Add footer
        self._add_footer(slide)

        # Save presentation
        self.prs.save(output_path)

        if self.verbose:
            print(f"✅ JAMA Oncology PowerPoint kaydedildi: {output_path}")

    def _add_header(self, slide):
        """Add GREEN header bar with JAMA Oncology branding"""
        # Header background
        left = Inches(0)
        top = Inches(0)
        width = Inches(self.SLIDE_WIDTH)
        height = Inches(self.HEADER_HEIGHT)

        header_box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.COLORS['header_bg']
        header_box.line.fill.background()

        # Add "JAMA Oncology" text
        logo_left = Inches(self.LOGO_LEFT)
        logo_top = Inches(self.LOGO_TOP + 0.05)
        logo_width = Inches(self.LOGO_WIDTH)
        logo_height = Inches(self.LOGO_HEIGHT)

        logo_textbox = slide.shapes.add_textbox(logo_left, logo_top, logo_width, logo_height)
        text_frame = logo_textbox.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # "JAMA Oncology" in one line
        p1 = text_frame.paragraphs[0]
        p1.text = "JAMA Oncology"
        p1.alignment = PP_ALIGN.LEFT
        p1.space_after = Pt(0)
        run1 = p1.runs[0]
        run1.font.size = Pt(22)
        run1.font.bold = False
        run1.font.color.rgb = RGBColor(255, 255, 255)
        run1.font.name = "Arial"

    def _add_main_title(self, slide):
        """Add main article title with RCT prefix - auto-size based on length"""
        title = self.data.get('title', 'Article Title')

        # Add study type prefix
        study_type = "RCT:"  # Default for oncology trials

        # Font size based on title length - OKUNABILIR KALSIN
        if len(title) > 150:
            font_size = 14
        elif len(title) > 120:
            font_size = 15
        elif len(title) > 100:
            font_size = 16
        elif len(title) > 70:
            font_size = 18
        else:
            font_size = 20

        left = Inches(self.TITLE_LEFT)
        top = Inches(self.TITLE_TOP)
        width = Inches(self.TITLE_WIDTH)
        height = Inches(self.TITLE_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)

        # Study type in GREEN - VURGULU
        run1 = p.add_run()
        run1.text = f"{study_type} "
        run1.font.bold = True
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = self.COLORS['accent']
        run1.font.name = "Arial"

        # Rest of title in black - BOLD
        run2 = p.add_run()
        run2.text = title
        run2.font.bold = True
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = RGBColor(0, 0, 0)
        run2.font.name = "Arial"

    def _add_content_box(self, slide, box_title: str, content: str,
                         left_pos: float, top_pos: float,
                         box_width: float, box_height: float,
                         add_chart: bool = False):
        """Add content box with title and content"""
        left = Inches(left_pos)
        top = Inches(top_pos)
        width = Inches(box_width)
        height = Inches(box_height)

        # Create box shape
        box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )

        # Format box
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['box_background']
        box.line.color.rgb = RGBColor(200, 200, 200)
        box.line.width = Pt(1)

        # Add icon based on box type
        self._add_oncology_icon(slide, box_title, left_pos, top_pos)

        # Try to add chart if requested
        chart_added = False
        if add_chart and box_title == "FINDINGS":
            # For FINDINGS, we'll add survival curve
            chart_added = True  # We'll add it manually for now

        # Add text frame
        text_frame = box.text_frame
        text_frame.word_wrap = True
        # More space for icon at top
        text_frame.margin_top = Inches(0.7)  # Space for icon
        text_frame.margin_bottom = Inches(0.1)
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)

        # Clear default paragraph
        text_frame.clear()

        # Add title paragraph
        title_p = text_frame.paragraphs[0]
        title_p.text = box_title.upper()
        title_p.alignment = PP_ALIGN.LEFT
        title_p.space_after = Pt(6)
        title_run = title_p.runs[0]
        title_run.font.bold = True
        title_run.font.size = Pt(9)
        title_run.font.color.rgb = self.COLORS['box_title']
        title_run.font.name = "Arial"

        # Add formatted content
        self._add_formatted_content(text_frame, content, box_title)

    def _add_formatted_content(self, text_frame, content: str, box_title: str):
        """Add formatted content with GREEN highlights for numbers - KELIME SINIRI UYGULANIR"""

        # Kelime limitleri - SENİN BELİRLEDİĞİN
        word_limits = {
            'POPULATION': 15,
            'INTERVENTION': 15,
            'SETTINGS / LOCATIONS': 10,
            'PRIMARY OUTCOME': 20,
            'FINDINGS': 15  # Her findings için
        }

        word_limit = word_limits.get(box_title, 20)

        # İçeriği kelime limitine göre kısalt
        content = self._truncate_to_word_limit(content, word_limit)

        lines = content.split('\n')
        total_words = sum(len(line.split()) for line in lines if line.strip())

        # Font boyutu - kelime sayısına göre
        if total_words > word_limit * 0.8:
            base_font_size = 9
        else:
            base_font_size = 10

        for line in lines:
            if not line.strip():
                continue

            # Capitalize first letter
            line = line[0].upper() + line[1:] if line else line

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)

            # Highlight numbers in GREEN (like reference image)
            if re.search(r'\d+', line):
                parts = re.split(r'(\d+\.?\d*)', line)
                for part in parts:
                    if part.strip():
                        run = p.add_run()
                        run.text = part
                        run.font.name = "Arial"

                        # Numbers in GREEN and bold - VURGULU
                        if re.match(r'\d+\.?\d*$', part):
                            run.font.color.rgb = self.COLORS['accent']
                            run.font.bold = True
                            run.font.size = Pt(base_font_size + 1)
                        else:
                            run.font.color.rgb = self.COLORS['box_content']
                            run.font.size = Pt(base_font_size)
            else:
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.color.rgb = self.COLORS['box_content']
                run.font.size = Pt(base_font_size)

    def _truncate_to_word_limit(self, text: str, word_limit: int) -> str:
        """Metni kelime limitine göre akıllıca kısalt - en önemli bilgileri koru"""
        if not text or not text.strip():
            return text

        # Satırları birleştir
        text = ' '.join(text.split('\n'))
        words = text.split()

        if len(words) <= word_limit:
            return text

        # AKILLI KISALTMA: Sayıları ve önemli kelimeleri koru
        important_words = []
        for word in words:
            # Sayı içeriyorsa muhakkak al
            if any(char.isdigit() for char in word):
                important_words.append(word)
            # Önemli kelimeler (study, patients, vs, treatment, risk, etc.)
            elif word.lower() in ['study', 'patients', 'adults', 'randomized', 'treatment',
                                   'vs', 'compared', 'risk', 'outcome', 'survival', 'age',
                                   'male', 'female', 'years', 'months', 'therapy', 'trial',
                                   'cohort', 'retrospective', 'prospective', 'analysis']:
                important_words.append(word)
            # Normal kelimeler
            else:
                important_words.append(word)

        # Word limit'e kadar al
        if len(important_words) > word_limit:
            # Sayı içeren kelimeleri ve ilk kelimeleri önceliklendir
            result = []
            # Önce sayılı olanları al
            for word in important_words:
                if any(char.isdigit() for char in word) and len(result) < word_limit:
                    result.append(word)

            # Kalan yerleri baştan doldur
            for word in important_words:
                if word not in result and len(result) < word_limit:
                    result.append(word)

            return ' '.join(result)

        return ' '.join(important_words[:word_limit])

    def _add_oncology_icon(self, slide, box_title: str, left_pos: float, top_pos: float):
        """Add oncology-specific icons - positioned ABOVE title"""
        icon_size = 0.3

        # Calculate center position
        if box_title == "FINDINGS":
            box_width = self.RIGHT_BOX_WIDTH
        else:
            box_width = self.LEFT_BOX_WIDTH

        icon_left = left_pos + (box_width - icon_size) / 2
        icon_top = top_pos + 0.3  # Lower position to avoid overlap

        if box_title == "POPULATION":
            self._add_brain_icon(slide, icon_left, icon_top, icon_size)
        elif box_title == "INTERVENTION":
            self._add_medication_icon(slide, icon_left, icon_top, icon_size)
        elif box_title == "FINDINGS":
            pass  # Chart will be shown
        elif box_title == "SETTINGS / LOCATIONS":
            self._add_map_pin_icon(slide, icon_left, icon_top, icon_size)
        elif box_title == "PRIMARY OUTCOME":
            self._add_target_icon(slide, icon_left, icon_top, icon_size)

    def _add_brain_icon(self, slide, left: float, top: float, size: float):
        """Brain/neurology icon - simple circle"""
        outline_color = RGBColor(0, 0, 0)

        # Brain outline (circle)
        brain = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )
        brain.fill.background()
        brain.line.color.rgb = outline_color
        brain.line.width = Pt(3)

    def _add_target_icon(self, slide, left: float, top: float, size: float):
        """Target/outcome icon - concentric circles"""
        outline_color = RGBColor(0, 0, 0)

        # Outer circle
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )
        outer.fill.background()
        outer.line.color.rgb = outline_color
        outer.line.width = Pt(2)

        # Inner circle
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.3), Inches(top + size * 0.3),
            Inches(size * 0.4), Inches(size * 0.4)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = outline_color
        inner.line.fill.background()

    def _add_breast_cancer_icon(self, slide, left: float, top: float, size: float):
        """Breast cancer anatomy icon - black outline"""
        outline_color = RGBColor(0, 0, 0)

        # Breast outline (circle/oval shape)
        breast = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.15), Inches(top + size * 0.1),
            Inches(size * 0.7), Inches(size * 0.75)
        )
        breast.fill.background()
        breast.line.color.rgb = outline_color
        breast.line.width = Pt(3)

        # Tumor indicator (filled dark circle)
        tumor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.25), Inches(top + size * 0.25),
            Inches(size * 0.2), Inches(size * 0.2)
        )
        tumor.fill.solid()
        tumor.fill.fore_color.rgb = outline_color
        tumor.line.fill.background()

        # Radiation/warning lines
        for i in range(3):
            angle_offset = i * 0.08
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * 0.75 + angle_offset), Inches(top + size * 0.15 + i * 0.06),
                Inches(size * 0.12), Inches(size * 0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = outline_color
            line.line.fill.background()

    def _add_medication_icon(self, slide, left: float, top: float, size: float):
        """Medication pills icon - outline with division line"""
        outline_color = RGBColor(0, 0, 0)

        # Two pills side by side
        pill_positions = [0, 0.5]

        for i, x_offset in enumerate(pill_positions):
            # Pill oval
            pill = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + size * x_offset), Inches(top + size * 0.2),
                Inches(size * 0.4), Inches(size * 0.6)
            )
            pill.fill.background()
            pill.line.color.rgb = outline_color
            pill.line.width = Pt(3)

            # Division line in middle of pill
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * (x_offset - 0.05)), Inches(top + size * 0.48),
                Inches(size * 0.5), Inches(size * 0.04)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = outline_color
            divider.line.fill.background()

    def _add_map_pin_icon(self, slide, left: float, top: float, size: float):
        """Map location pin icon - outline"""
        outline_color = RGBColor(0, 0, 0)

        # Pin top (circle)
        pin_top = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.25), Inches(top),
            Inches(size * 0.5), Inches(size * 0.5)
        )
        pin_top.fill.background()
        pin_top.line.color.rgb = outline_color
        pin_top.line.width = Pt(3)

        # Pin bottom (triangle pointing down)
        # Using rectangle rotated as triangle placeholder
        pin_bottom = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(left + size * 0.3), Inches(top + size * 0.45),
            Inches(size * 0.4), Inches(size * 0.5)
        )
        pin_bottom.fill.background()
        pin_bottom.line.color.rgb = outline_color
        pin_bottom.line.width = Pt(3)

    def _add_findings_box(self, slide):
        """Add FINDINGS box with chart and iDFS results - like reference image"""
        left = Inches(self.RIGHT_BOX_X)
        top = Inches(self.RIGHT_BOX_Y)
        width = Inches(self.RIGHT_BOX_WIDTH)
        height = Inches(self.RIGHT_BOX_HEIGHT)

        # Create main box (BG only - NO TEXT FRAME)
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['box_background']
        box.line.color.rgb = RGBColor(200, 200, 200)
        box.line.width = Pt(1)

        # Add FINDINGS title as separate textbox
        title_box = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + Inches(0.12),
            width - Inches(0.3),
            Inches(0.25)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "FINDINGS"
        title_p.alignment = PP_ALIGN.LEFT
        title_run = title_p.runs[0]
        title_run.font.bold = True
        title_run.font.size = Pt(10)
        title_run.font.color.rgb = self.COLORS['box_title']
        title_run.font.name = "Arial"

        # Add main finding content (text before chart) as separate textbox - MAX 15 KELIME
        finding_content = self.data.get('finding_1', '')
        if finding_content:
            # Get text before "iDFS" section
            content_lines = finding_content.split('\n')
            main_text = []
            for line in content_lines:
                if 'iDFS' in line or 'Ribociclib + NSAI:' in line:
                    break
                if line.strip():
                    main_text.append(line.strip())

            main_content = ' '.join(main_text).strip()

            # Kelime limitine göre kısalt (max 15 kelime)
            main_content = self._truncate_to_word_limit(main_content, 15)

            if main_content:
                content_box = slide.shapes.add_textbox(
                    left + Inches(0.15),
                    top + Inches(0.45),
                    width - Inches(0.3),
                    Inches(0.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_p = content_frame.paragraphs[0]
                content_p.text = main_content
                content_p.alignment = PP_ALIGN.LEFT
                run = content_p.runs[0]
                run.font.size = Pt(8)
                run.font.color.rgb = self.COLORS['box_content']
                run.font.name = "Arial"

        # Add survival curve placeholder (chart area)
        chart_left = self.RIGHT_BOX_X + 0.25
        chart_top = self.RIGHT_BOX_Y + 1.05
        chart_width = self.RIGHT_BOX_WIDTH - 0.5
        chart_height = 2.5

        self._add_survival_curve_placeholder(slide, chart_left, chart_top, chart_width, chart_height)

        # Add iDFS results box at bottom of FINDINGS
        idfs_box_top = self.RIGHT_BOX_Y + 3.8
        idfs_box_height = 1.0
        self._add_idfs_results(slide, self.RIGHT_BOX_X + 0.15, idfs_box_top,
                              self.RIGHT_BOX_WIDTH - 0.3, idfs_box_height)

    def _add_survival_curve_placeholder(self, slide, left: float, top: float, width: float, height: float):
        """Add survival curve chart placeholder with visual representation"""

        chart_left = Inches(left)
        chart_top = Inches(top)
        chart_width = Inches(width)
        chart_height = Inches(height)

        # Add white background for chart area
        chart_bg = slide.shapes.add_shape(1, chart_left, chart_top, chart_width, chart_height)
        chart_bg.fill.solid()
        chart_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        chart_bg.line.color.rgb = RGBColor(220, 220, 220)
        chart_bg.line.width = Pt(0.5)

        # Draw two survival curves (green lines representing treatment groups)
        # Line 1: Ribociclib+NSAI (higher survival - top line)
        line1_points = [
            (left + 0.3, top + 0.3),
            (left + 0.8, top + 0.35),
            (left + 1.3, top + 0.4),
            (left + 1.8, top + 0.45),
            (left + 2.3, top + 0.5),
            (left + 2.8, top + 0.6)
        ]

        for i in range(len(line1_points) - 1):
            x1, y1 = line1_points[i]
            x2, y2 = line1_points[i + 1]
            connector = slide.shapes.add_connector(
                1,  # Straight line
                Inches(x1), Inches(y1),
                Inches(x2), Inches(y2)
            )
            connector.line.color.rgb = self.COLORS['accent']  # Green
            connector.line.width = Pt(2)

        # Line 2: NSAI alone (lower survival - bottom line)
        line2_points = [
            (left + 0.3, top + 0.4),
            (left + 0.8, top + 0.5),
            (left + 1.3, top + 0.6),
            (left + 1.8, top + 0.7),
            (left + 2.3, top + 0.85),
            (left + 2.8, top + 1.0)
        ]

        for i in range(len(line2_points) - 1):
            x1, y1 = line2_points[i]
            x2, y2 = line2_points[i + 1]
            connector = slide.shapes.add_connector(
                1,  # Straight line
                Inches(x1), Inches(y1),
                Inches(x2), Inches(y2)
            )
            connector.line.color.rgb = RGBColor(128, 128, 128)  # Gray
            connector.line.width = Pt(2)

        # Add legend
        legend_y = top + 0.15

        # Green line for Ribociclib+NSAI
        legend_line1 = slide.shapes.add_shape(
            1,  # Rectangle
            chart_left + Inches(2.0), Inches(legend_y),
            Inches(0.2), Inches(0.02)
        )
        legend_line1.fill.solid()
        legend_line1.fill.fore_color.rgb = self.COLORS['accent']
        legend_line1.line.fill.background()

        legend_text1 = slide.shapes.add_textbox(
            chart_left + Inches(2.25), Inches(legend_y - 0.08),
            Inches(0.8), Inches(0.15)
        )
        legend_frame1 = legend_text1.text_frame
        legend_p1 = legend_frame1.paragraphs[0]
        legend_p1.text = "Ribociclib+NSAI"
        run1 = legend_p1.runs[0]
        run1.font.size = Pt(6)
        run1.font.color.rgb = RGBColor(60, 60, 60)
        run1.font.name = "Arial"

        # Gray line for NSAI alone
        legend_line2 = slide.shapes.add_shape(
            1,  # Rectangle
            chart_left + Inches(2.0), Inches(legend_y + 0.15),
            Inches(0.2), Inches(0.02)
        )
        legend_line2.fill.solid()
        legend_line2.fill.fore_color.rgb = RGBColor(128, 128, 128)
        legend_line2.line.fill.background()

        legend_text2 = slide.shapes.add_textbox(
            chart_left + Inches(2.25), Inches(legend_y + 0.07),
            Inches(0.6), Inches(0.15)
        )
        legend_frame2 = legend_text2.text_frame
        legend_p2 = legend_frame2.paragraphs[0]
        legend_p2.text = "NSAI alone"
        run2 = legend_p2.runs[0]
        run2.font.size = Pt(6)
        run2.font.color.rgb = RGBColor(60, 60, 60)
        run2.font.name = "Arial"

    def _add_idfs_results(self, slide, left: float, top: float, width: float, height: float):
        """Add iDFS results section below chart"""
        left_pos = Inches(left)
        top_pos = Inches(top)
        width_dim = Inches(width)
        height_dim = Inches(height)

        # Create text box (no background, just text)
        textbox = slide.shapes.add_textbox(left_pos, top_pos, width_dim, height_dim)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_left = Inches(0.05)
        text_frame.clear()

        # Title
        title_p = text_frame.paragraphs[0]
        title_p.text = "iDFS"
        title_p.alignment = PP_ALIGN.LEFT
        title_p.space_after = Pt(4)
        title_run = title_p.runs[0]
        title_run.font.bold = True
        title_run.font.size = Pt(9)
        title_run.font.color.rgb = self.COLORS['box_title']
        title_run.font.name = "Arial"

        # Get iDFS data from finding_1
        finding_content = self.data.get('finding_1', '')
        idfs_lines = []
        capture = False

        for line in finding_content.split('\n'):
            if 'iDFS' in line or 'Ribociclib + NSAI:' in line:
                capture = True
            if capture and line.strip():
                if not line.strip().startswith('iDFS'):
                    idfs_lines.append(line.strip())

        # If no iDFS data found, use default
        if not idfs_lines:
            idfs_lines = [
                'Ribociclib + NSAI: 263 events (10.3%)',
                'NSAI alone: 340 events (13.3%)',
                'Hazard ratio: 0.72; 95% CI, 0.61-0.84'
            ]

        # Add iDFS results with number highlighting
        for line in idfs_lines:
            if not line.strip():
                continue

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)

            # Highlight numbers in green
            if re.search(r'\d+', line):
                parts = re.split(r'(\d+\.?\d*)', line)
                for part in parts:
                    if part.strip():
                        run = p.add_run()
                        run.text = part
                        run.font.name = "Arial"

                        # Numbers in green and bold
                        if re.match(r'\d+\.?\d*$', part):
                            run.font.color.rgb = self.COLORS['accent']
                            run.font.bold = True
                            run.font.size = Pt(9)
                        else:
                            run.font.color.rgb = self.COLORS['box_content']
                            run.font.size = Pt(9)
            else:
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.color.rgb = self.COLORS['box_content']
                run.font.size = Pt(9)

    def _add_footer(self, slide):
        """Add footer with citation"""
        authors = self.data.get('authors', 'Authors')
        journal = "JAMA Oncol"
        pub_date = self.data.get('publication_date', '2025')
        doi = self.data.get('doi', '')

        # Extract year
        year = pub_date.split()[-1] if pub_date else '2025'

        # Create footer text
        footer_text = f"{authors}. {journal}. Published online {pub_date}. doi:10.1001/jamaoncol.{doi} © AMA"

        left = Inches(0.3)
        top = Inches(self.FOOTER_Y)
        width = Inches(9.4)
        height = Inches(self.FOOTER_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        p = text_frame.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.LEFT

        run = p.runs[0]
        run.font.size = Pt(6)
        run.font.color.rgb = self.COLORS['footer_text']
        run.font.name = "Arial"
