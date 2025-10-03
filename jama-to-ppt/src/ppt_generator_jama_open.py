"""
JAMA Network Open Format PowerPoint Generator
"""

from typing import Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import re

try:
    from chart_generator import ChartGenerator
    CHART_AVAILABLE = True
except ImportError:
    CHART_AVAILABLE = False


class JAMAOpenPowerPointGenerator:
    """
    Generates PowerPoint presentations in JAMA Network Open format
    """

    # JAMA Open Color Palette
    COLORS = {
        'header_bg': RGBColor(233, 19, 131),         # #E91383 - Magenta/Pink
        'title_text': RGBColor(255, 255, 255),       # #FFFFFF - White
        'box_background': RGBColor(240, 240, 240),   # #F0F0F0 - Light gray
        'box_title': RGBColor(153, 0, 102),          # #990066 - Dark magenta
        'box_content': RGBColor(50, 50, 50),         # #323232 - Dark gray
        'footer_text': RGBColor(100, 100, 100),      # #646464 - Gray
        'accent': RGBColor(233, 19, 131),            # #E91383 - Magenta
    }

    # Dimensions (in inches)
    SLIDE_WIDTH = 10
    SLIDE_HEIGHT = 7.5

    # Header
    HEADER_HEIGHT = 0.8
    LOGO_WIDTH = 2.5
    LOGO_HEIGHT = 0.6
    LOGO_LEFT = 0.3
    LOGO_TOP = 0.15

    # Main title
    TITLE_TOP = 1.0
    TITLE_LEFT = 0.4
    TITLE_WIDTH = 9.2
    TITLE_HEIGHT = 0.7

    # Content boxes - Top row (3 columns)
    TOP_BOX_WIDTH = 2.9
    TOP_BOX_HEIGHT = 2.8
    TOP_ROW_Y = 1.9

    POPULATION_X = 0.4
    INTERVENTION_X = 3.5
    FINDINGS_X = 6.6

    # Content boxes - Bottom row (2 columns)
    BOTTOM_BOX_WIDTH = 4.4
    BOTTOM_BOX_HEIGHT = 2.0
    BOTTOM_ROW_Y = 4.9

    SETTINGS_X = 0.4
    PRIMARY_OUTCOME_X = 5.1

    # Footer
    FOOTER_Y = 7.1
    FOOTER_HEIGHT = 0.3

    def __init__(self, article_data: Dict[str, str], icon_type: str = 'medical', verbose: bool = False):
        self.data = article_data
        self.icon_type = icon_type
        self.verbose = verbose
        self.prs = Presentation()

        # Set slide dimensions
        self.prs.slide_width = Inches(self.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT)

    def generate(self, output_path: str):
        """Generate PowerPoint presentation"""
        if self.verbose:
            print("📝 JAMA Open PowerPoint oluşturuluyor...")

        # Add blank slide
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Add header bar
        self._add_header(slide)

        # Add main title
        self._add_main_title(slide)

        # Add top row boxes (3 columns)
        self._add_content_box(
            slide, "POPULATION", self.data.get('population', 'N/A'),
            self.POPULATION_X, self.TOP_ROW_Y, self.TOP_BOX_WIDTH, self.TOP_BOX_HEIGHT
        )
        self._add_content_box(
            slide, "INTERVENTION", self.data.get('intervention', 'N/A'),
            self.INTERVENTION_X, self.TOP_ROW_Y, self.TOP_BOX_WIDTH, self.TOP_BOX_HEIGHT
        )
        # Add FINDINGS box with optional chart
        findings_content = self.data.get('finding_1', 'N/A')
        self._add_content_box(
            slide, "FINDINGS", findings_content,
            self.FINDINGS_X, self.TOP_ROW_Y, self.TOP_BOX_WIDTH, self.TOP_BOX_HEIGHT,
            add_chart=True
        )

        # Add bottom row boxes (2 columns)
        self._add_content_box(
            slide, "SETTINGS / LOCATIONS", self.data.get('setting', 'N/A'),
            self.SETTINGS_X, self.BOTTOM_ROW_Y, self.BOTTOM_BOX_WIDTH, self.BOTTOM_BOX_HEIGHT
        )
        self._add_content_box(
            slide, "PRIMARY OUTCOME", self.data.get('primary_outcome', 'N/A'),
            self.PRIMARY_OUTCOME_X, self.BOTTOM_ROW_Y, self.BOTTOM_BOX_WIDTH, self.BOTTOM_BOX_HEIGHT
        )

        # Add footer
        self._add_footer(slide)

        # Save presentation
        self.prs.save(output_path)

        if self.verbose:
            print(f"✅ JAMA Open PowerPoint kaydedildi: {output_path}")

    def _add_header(self, slide):
        """Add pink header bar with JAMA Network Open branding"""
        # Header background
        left = Inches(0)
        top = Inches(0)
        width = Inches(self.SLIDE_WIDTH)
        height = Inches(self.HEADER_HEIGHT)

        header_box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.COLORS['header_bg']
        header_box.line.fill.background()

        # Add "JAMA Network Open" text logo
        logo_left = Inches(self.LOGO_LEFT)
        logo_top = Inches(self.LOGO_TOP)
        logo_width = Inches(self.LOGO_WIDTH)
        logo_height = Inches(self.LOGO_HEIGHT)

        logo_textbox = slide.shapes.add_textbox(logo_left, logo_top, logo_width, logo_height)
        text_frame = logo_textbox.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # "JAMA" line
        p1 = text_frame.paragraphs[0]
        p1.text = "JAMA"
        p1.alignment = PP_ALIGN.LEFT
        p1.space_after = Pt(0)
        run1 = p1.runs[0]
        run1.font.size = Pt(16)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(255, 255, 255)
        run1.font.name = "Arial"

        # "Network" line
        p2 = text_frame.add_paragraph()
        p2.text = "Network "
        p2.alignment = PP_ALIGN.LEFT
        p2.space_after = Pt(0)
        run2 = p2.runs[0]
        run2.font.size = Pt(10)
        run2.font.color.rgb = RGBColor(255, 255, 255)
        run2.font.name = "Arial"

        # Add "Open" in same line
        run3 = p2.add_run()
        run3.text = "Open"
        run3.font.size = Pt(16)
        run3.font.bold = True
        run3.font.color.rgb = RGBColor(255, 255, 255)
        run3.font.name = "Arial"

    def _add_main_title(self, slide):
        """Add main article title"""
        title = self.data.get('title', 'Article Title')

        # Add study type prefix if available
        study_type = self._detect_study_type(title)
        if study_type:
            title = f"{study_type}: {title}"

        # Determine font size based on title length
        if len(title) > 100:
            font_size = 18
        elif len(title) > 70:
            font_size = 20
        else:
            font_size = 24

        left = Inches(self.TITLE_LEFT)
        top = Inches(self.TITLE_TOP)
        width = Inches(self.TITLE_WIDTH)
        height = Inches(self.TITLE_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Add study type in color
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)

        if study_type:
            # Study type in magenta
            run1 = p.add_run()
            run1.text = f"{study_type}: "
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = self.COLORS['accent']
            run1.font.name = "Arial"

            # Rest of title in black
            run2 = p.add_run()
            run2.text = title.replace(f"{study_type}: ", "")
            run2.font.bold = True
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = RGBColor(0, 0, 0)
            run2.font.name = "Arial"
        else:
            p.text = title
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.name = "Arial"

    def _detect_study_type(self, title: str) -> str:
        """Detect study type from title or data"""
        title_lower = title.lower()

        # Check common study types
        study_types = {
            'RCT': ['randomized', 'rct', 'clinical trial'],
            'Cohort Study': ['cohort'],
            'Meta-Analysis': ['meta-analysis', 'systematic review'],
            'Case-Control': ['case-control'],
            'Cross-Sectional': ['cross-sectional'],
        }

        for study_type, keywords in study_types.items():
            for keyword in keywords:
                if keyword in title_lower:
                    return study_type

        return ''

    def _add_content_box(self, slide, box_title: str, content: str,
                         left_pos: float, top_pos: float,
                         box_width: float, box_height: float,
                         add_chart: bool = False):
        """Add content box with title and content"""
        left = Inches(left_pos)
        top = Inches(top_pos)
        width = Inches(box_width)
        height = Inches(box_height)

        # Create box shape
        box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )

        # Format box
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['box_background']
        box.line.color.rgb = RGBColor(200, 200, 200)
        box.line.width = Pt(1)

        # Add icon based on box type
        self._add_simple_icon(slide, box_title, left_pos, top_pos)

        # Try to add chart if requested and available
        chart_added = False
        if add_chart and CHART_AVAILABLE and box_title == "FINDINGS":
            chart_added = self._try_add_chart(slide, content, left_pos, top_pos, box_width, box_height)

        # Add text frame
        text_frame = box.text_frame
        text_frame.word_wrap = True
        # Compact margins - less space
        text_frame.margin_top = Inches(0.5)  # Reduced from 0.55
        text_frame.margin_bottom = Inches(0.1)  # Reduced from 0.15
        text_frame.margin_left = Inches(0.12)  # Reduced from 0.15
        text_frame.margin_right = Inches(0.12)  # Reduced from 0.15

        # Clear default paragraph
        text_frame.clear()

        # Add title paragraph
        title_p = text_frame.paragraphs[0]
        title_p.text = box_title.upper()
        title_p.alignment = PP_ALIGN.LEFT
        title_p.space_after = Pt(4)  # Reduced from 8 - less space after title
        title_run = title_p.runs[0]
        title_run.font.bold = True
        title_run.font.size = Pt(10)  # Reduced from 11 - slightly smaller
        title_run.font.color.rgb = self.COLORS['box_title']
        title_run.font.name = "Arial"

        # Add formatted content with proper styling
        self._add_formatted_content(text_frame, content, box_title)

    def _add_formatted_content(self, text_frame, content: str, box_title: str):
        """Add formatted content with colored numbers and bold text"""
        # Adjust font size based on content length
        if len(content) > 400:
            base_font_size = 9
        elif len(content) > 250:
            base_font_size = 10
        else:
            base_font_size = 11

        # Box-specific formatting
        if box_title == "POPULATION":
            self._format_population_content(text_frame, content, base_font_size)
        elif box_title == "INTERVENTION":
            self._format_intervention_content(text_frame, content, base_font_size)
        elif box_title == "FINDINGS":
            self._format_findings_content(text_frame, content, base_font_size)
        elif box_title == "SETTINGS / LOCATIONS":
            self._format_settings_content(text_frame, content, base_font_size)
        elif box_title == "PRIMARY OUTCOME":
            self._format_outcome_content(text_frame, content, base_font_size)
        else:
            # Default formatting
            p = text_frame.add_paragraph()
            p.text = content
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(base_font_size)
            run.font.color.rgb = self.COLORS['box_content']
            run.font.name = "Arial"

    def _format_population_content(self, text_frame, content: str, font_size: int):
        """Format POPULATION content with bold numbers"""
        lines = content.split('\n')

        for line in lines:
            if not line.strip():
                continue

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)  # Reduced from 4 - tighter spacing

            # Check for patterns like "130 Men, 14 Women" or "Mean age, 67 y"
            if re.search(r'\d+', line):
                # Bold numbers and key info
                parts = re.split(r'(\d+)', line)
                for part in parts:
                    if part.strip():
                        run = p.add_run()
                        run.text = part
                        run.font.name = "Arial"
                        run.font.color.rgb = self.COLORS['box_content']

                        if part.isdigit() or 'Mean age' in line or 'Men' in part or 'Women' in part:
                            run.font.bold = True
                            run.font.size = Pt(font_size + 1)
                        else:
                            run.font.size = Pt(font_size)
            else:
                run = p.add_run()
                run.text = line
                run.font.size = Pt(font_size)
                run.font.color.rgb = self.COLORS['box_content']
                run.font.name = "Arial"

    def _format_intervention_content(self, text_frame, content: str, font_size: int):
        """Format INTERVENTION content with colored numbers"""
        lines = content.split('\n')

        for line in lines:
            if not line.strip():
                continue

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)  # Reduced from 4 - tighter spacing

            # Highlight numbers in magenta and make bold
            if re.search(r'\d+', line):
                parts = re.split(r'(\d+\s*Sites|\d+\s*patients|\d+)', line)
                for part in parts:
                    if part.strip():
                        run = p.add_run()
                        run.text = part
                        run.font.name = "Arial"

                        # Numbers and "Sites" in magenta and bold
                        if re.search(r'\d+', part) or 'Sites' in part:
                            run.font.color.rgb = self.COLORS['accent']
                            run.font.bold = True
                            run.font.size = Pt(font_size + 1)
                        # Section headers bold
                        elif 'Enhanced support' in part or 'Foundational support' in part:
                            run.font.color.rgb = self.COLORS['box_content']
                            run.font.bold = True
                            run.font.size = Pt(font_size)
                        else:
                            run.font.color.rgb = self.COLORS['box_content']
                            run.font.size = Pt(font_size)
            else:
                # Headers like "Enhanced support"
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.color.rgb = self.COLORS['box_content']

                if 'support' in line.lower() or 'Sites' in line:
                    run.font.bold = True

                run.font.size = Pt(font_size)

    def _format_findings_content(self, text_frame, content: str, font_size: int):
        """Format FINDINGS content"""
        lines = content.split('\n')

        for line in lines:
            if not line.strip():
                continue

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)  # Reduced from 4 - tighter spacing

            # "Mean difference" in magenta
            if 'Mean difference' in line or 'mean difference' in line:
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.color.rgb = self.COLORS['accent']
                run.font.bold = True
                run.font.size = Pt(font_size)
            else:
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.color.rgb = self.COLORS['box_content']
                run.font.size = Pt(font_size - 1)

    def _format_settings_content(self, text_frame, content: str, font_size: int):
        """Format SETTINGS content with VA emphasis"""
        lines = content.split('\n')

        for line in lines:
            if not line.strip():
                continue

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)  # Reduced from 4 - tighter spacing

            # Highlight "VA" and numbers
            parts = re.split(r'(\d+\s*VA|VA|\d+)', line)
            for part in parts:
                if part.strip():
                    run = p.add_run()
                    run.text = part
                    run.font.name = "Arial"

                    if 'VA' in part or part.isdigit():
                        run.font.bold = True
                        run.font.size = Pt(font_size + 2)
                        run.font.color.rgb = self.COLORS['box_content']
                    else:
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = self.COLORS['box_content']

    def _format_outcome_content(self, text_frame, content: str, font_size: int):
        """Format PRIMARY OUTCOME content"""
        p = text_frame.add_paragraph()
        p.text = content
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(font_size - 1)
        run.font.color.rgb = self.COLORS['box_content']
        run.font.name = "Arial"

    def _add_footer(self, slide):
        """Add footer with citation info"""
        # Format: Authors et al. Journal. Year;Volume(Issue):Pages. doi:DOI
        authors = self.data.get('authors', 'Authors')
        journal = "JAMA Netw Open"
        pub_date = self.data.get('publication_date', '2025')
        doi = self.data.get('doi', '')

        # Extract year from publication date
        year = pub_date.split()[-1] if pub_date else '2025'

        # Create footer text
        footer_text = f"{authors}. {journal}. {year};{doi}. doi:10.1001/jamanetworkopen.{doi}"

        # Add copyright
        footer_text += " © AMA"

        left = Inches(0.4)
        top = Inches(self.FOOTER_Y)
        width = Inches(9.2)
        height = Inches(self.FOOTER_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.LEFT

        run = p.runs[0]
        run.font.size = Pt(8)
        run.font.color.rgb = self.COLORS['footer_text']
        run.font.name = "Arial"

    def _add_simple_icon(self, slide, box_title: str, left_pos: float, top_pos: float):
        """Add simple magenta icon at top of box - like reference image"""
        icon_size = 0.35  # Reduced from 0.4 - smaller icons
        # Center icon horizontally at top of box

        # Map box titles to icon types
        icon_map = {
            'POPULATION': self._create_knee_icon,
            'INTERVENTION': self._create_people_group_icon,
            'FINDINGS': self._create_chart_icon,
            'SETTINGS / LOCATIONS': self._create_building_icon,
            'PRIMARY OUTCOME': self._create_target_icon,
        }

        # Add icon if mapping exists
        if box_title in icon_map:
            # Calculate center position
            if box_title in ['POPULATION', 'INTERVENTION', 'FINDINGS']:
                box_width = self.TOP_BOX_WIDTH
            else:
                box_width = self.BOTTOM_BOX_WIDTH

            icon_left = left_pos + (box_width - icon_size) / 2
            icon_top = top_pos + 0.12  # Reduced from 0.15 - closer to top

            icon_map[box_title](slide, icon_left, icon_top, icon_size)

    def _create_knee_icon(self, slide, left: float, top: float, size: float):
        """Simple knee icon - like in reference image"""
        icon_color = self.COLORS['accent']

        # Simple knee outline - curved shapes
        # Upper leg
        upper = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + size * 0.2), Inches(top),
            Inches(size * 0.25), Inches(size * 0.45)
        )
        upper.fill.solid()
        upper.fill.fore_color.rgb = icon_color
        upper.line.fill.background()
        upper.rotation = -15

        # Knee joint (oval)
        knee = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.3), Inches(top + size * 0.35),
            Inches(size * 0.3), Inches(size * 0.35)
        )
        knee.fill.solid()
        knee.fill.fore_color.rgb = icon_color
        knee.line.fill.background()

        # Lower leg
        lower = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + size * 0.35), Inches(top + size * 0.55),
            Inches(size * 0.22), Inches(size * 0.4)
        )
        lower.fill.solid()
        lower.fill.fore_color.rgb = icon_color
        lower.line.fill.background()
        lower.rotation = 10

        # Pain lines
        for i in range(3):
            pain = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * 0.65 + i * 0.04), Inches(top + size * 0.4),
                Inches(size * 0.12), Inches(size * 0.02)
            )
            pain.fill.solid()
            pain.fill.fore_color.rgb = icon_color
            pain.line.fill.background()

    def _add_population_icon(self, slide, left: float, top: float, size: float):
        """Add knee illustration for population"""
        # Draw a simplified knee joint illustration
        # Using black color for medical illustration style
        knee_color = RGBColor(0, 0, 0)

        # Femur (upper leg bone) - slightly angled
        femur = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + size * 0.15), Inches(top),
            Inches(size * 0.25), Inches(size * 0.4)
        )
        femur.fill.solid()
        femur.fill.fore_color.rgb = RGBColor(255, 255, 255)
        femur.line.color.rgb = knee_color
        femur.line.width = Pt(2.5)
        femur.rotation = -10

        # Knee cap (patella) - oval
        patella = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.3), Inches(top + size * 0.35),
            Inches(size * 0.25), Inches(size * 0.3)
        )
        patella.fill.solid()
        patella.fill.fore_color.rgb = RGBColor(255, 255, 255)
        patella.line.color.rgb = knee_color
        patella.line.width = Pt(2.5)

        # Tibia (lower leg bone) - slightly angled opposite direction
        tibia = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + size * 0.25), Inches(top + size * 0.55),
            Inches(size * 0.22), Inches(size * 0.4)
        )
        tibia.fill.solid()
        tibia.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tibia.line.color.rgb = knee_color
        tibia.line.width = Pt(2.5)
        tibia.rotation = 5

        # Pain indicator lines (like in medical diagrams)
        for i in range(3):
            pain_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * 0.65 + i * 0.03), Inches(top + size * 0.4 + i * 0.02),
                Inches(size * 0.15), Inches(size * 0.02)
            )
            pain_line.fill.solid()
            pain_line.fill.fore_color.rgb = knee_color
            pain_line.line.fill.background()

    def _create_people_group_icon(self, slide, left: float, top: float, size: float):
        """Group of people icon - like reference image"""
        icon_color = self.COLORS['accent']

        # Create 3 person silhouettes side by side
        positions = [0, 0.35, 0.7]

        for i, x_offset in enumerate(positions):
            # Make middle person slightly taller
            height_mult = 0.55 if i == 1 else 0.5
            y_offset = 0 if i == 1 else 0.05

            # Head
            head = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + size * x_offset), Inches(top + size * y_offset),
                Inches(size * 0.25), Inches(size * 0.25)
            )
            head.fill.solid()
            head.fill.fore_color.rgb = icon_color
            head.line.fill.background()

            # Body
            body = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + size * (x_offset - 0.05)), Inches(top + size * (0.25 + y_offset)),
                Inches(size * 0.35), Inches(size * height_mult)
            )
            body.fill.solid()
            body.fill.fore_color.rgb = icon_color
            body.line.fill.background()

    def _add_intervention_icon(self, slide, left: float, top: float, size: float):
        """Add intervention icon (medical cross or syringe)"""
        icon_color = self.COLORS['accent']

        # Medical cross
        # Vertical bar
        v_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + size * 0.35), Inches(top),
            Inches(size * 0.3), Inches(size)
        )
        v_bar.fill.solid()
        v_bar.fill.fore_color.rgb = icon_color
        v_bar.line.fill.background()

        # Horizontal bar
        h_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top + size * 0.35),
            Inches(size), Inches(size * 0.3)
        )
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = icon_color
        h_bar.line.fill.background()

    def _create_chart_icon(self, slide, left: float, top: float, size: float):
        """Bar chart icon - like reference image"""
        icon_color = self.COLORS['accent']

        # Three bars
        heights = [0.5, 0.8, 0.6]
        for i, h in enumerate(heights):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + i * size * 0.35), Inches(top + size * (1 - h)),
                Inches(size * 0.25), Inches(size * h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = icon_color
            bar.line.fill.background()

    def _add_chart_icon(self, slide, left: float, top: float, size: float):
        """Add chart icon (bar chart representation)"""
        icon_color = self.COLORS['accent']

        # Three bars of different heights
        heights = [0.4, 0.7, 0.5]
        for i, h in enumerate(heights):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + i * size * 0.35), Inches(top + size * (1 - h)),
                Inches(size * 0.25), Inches(size * h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = icon_color
            bar.line.fill.background()

    def _create_building_icon(self, slide, left: float, top: float, size: float):
        """Simple building icon - like reference image"""
        icon_color = self.COLORS['accent']

        # Main building
        building = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + size * 0.1), Inches(top + size * 0.25),
            Inches(size * 0.8), Inches(size * 0.7)
        )
        building.fill.solid()
        building.fill.fore_color.rgb = icon_color
        building.line.fill.background()

        # Windows (white rectangles)
        for row in range(2):
            for col in range(3):
                window = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left + size * (0.2 + col * 0.25)), Inches(top + size * (0.35 + row * 0.25)),
                    Inches(size * 0.12), Inches(size * 0.15)
                )
                window.fill.solid()
                window.fill.fore_color.rgb = RGBColor(255, 255, 255)
                window.line.fill.background()

        # Door
        door = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + size * 0.4), Inches(top + size * 0.75),
            Inches(size * 0.2), Inches(size * 0.2)
        )
        door.fill.solid()
        door.fill.fore_color.rgb = RGBColor(255, 255, 255)
        door.line.fill.background()

    def _create_target_icon(self, slide, left: float, top: float, size: float):
        """Target/bullseye icon - like reference image"""
        icon_color = self.COLORS['accent']

        # Outer circle
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.15), Inches(top + size * 0.15),
            Inches(size * 0.7), Inches(size * 0.7)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = icon_color
        outer.line.fill.background()

        # Inner white circle
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.3), Inches(top + size * 0.3),
            Inches(size * 0.4), Inches(size * 0.4)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = RGBColor(255, 255, 255)
        inner.line.fill.background()

        # Center dot
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.42), Inches(top + size * 0.42),
            Inches(size * 0.16), Inches(size * 0.16)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = icon_color
        center.line.fill.background()

    def _add_location_icon(self, slide, left: float, top: float, size: float):
        """Add VA building icon - hospital with VA label"""
        icon_color = RGBColor(0, 0, 0)

        # Main building outline (simple rectangle)
        building = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top + size * 0.3),
            Inches(size * 0.9), Inches(size * 0.6)
        )
        building.fill.solid()
        building.fill.fore_color.rgb = RGBColor(255, 255, 255)
        building.line.color.rgb = icon_color
        building.line.width = Pt(2.5)

        # VA text box inside building
        va_box = slide.shapes.add_textbox(
            Inches(left + size * 0.1), Inches(top + size * 0.4),
            Inches(size * 0.7), Inches(size * 0.4)
        )
        va_frame = va_box.text_frame
        va_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        va_p = va_frame.paragraphs[0]
        va_p.text = "VA"
        va_p.alignment = PP_ALIGN.CENTER

        va_run = va_p.runs[0]
        va_run.font.size = Pt(18)
        va_run.font.bold = True
        va_run.font.color.rgb = icon_color
        va_run.font.name = "Arial"

        # Building entrance (small rectangle at bottom)
        door = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + size * 0.35), Inches(top + size * 0.75),
            Inches(size * 0.2), Inches(size * 0.15)
        )
        door.fill.solid()
        door.fill.fore_color.rgb = icon_color
        door.line.fill.background()

    def _add_outcome_icon(self, slide, left: float, top: float, size: float):
        """Add outcome icon (target/bullseye)"""
        icon_color = self.COLORS['accent']

        # Outer circle
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = icon_color
        outer.line.fill.background()

        # Inner circle (white)
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.25), Inches(top + size * 0.25),
            Inches(size * 0.5), Inches(size * 0.5)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = RGBColor(255, 255, 255)
        inner.line.fill.background()

        # Center dot
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.4), Inches(top + size * 0.4),
            Inches(size * 0.2), Inches(size * 0.2)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = icon_color
        center.line.fill.background()

    def _extract_chart_data(self, content: str) -> Optional[Dict]:
        """Extract numerical data from findings for chart generation"""
        # Look for patterns like "Enhanced: 1.0, Foundational: 1.0"
        # or "P=.92" or percentages

        patterns = {
            'mean': r'[Mm]ean.*?([0-9.]+)',
            'p_value': r'[Pp]\s*=\s*([0-9.]+)',
            'percentage': r'(\d+(?:\.\d+)?)\s*%',
            'ci': r'95%\s*CI[,:]?\s*([0-9.-]+)\s+to\s+([0-9.-]+)',
        }

        data = {}
        for key, pattern in patterns.items():
            matches = re.findall(pattern, content)
            if matches:
                data[key] = matches

        return data if data else None

    def _try_add_chart(self, slide, content: str, left_pos: float, top_pos: float,
                      box_width: float, box_height: float) -> bool:
        """Try to add a chart based on content. Returns True if successful."""
        if not CHART_AVAILABLE:
            return False

        try:
            # Try to generate chart from content
            chart_buffer = ChartGenerator.create_comparison_chart(content, width=800, height=600)

            if chart_buffer:
                # Add chart image to slide
                # Position in lower part of findings box
                chart_left = Inches(left_pos + 0.3)
                chart_top = Inches(top_pos + 1.2)
                chart_width = Inches(box_width - 0.6)
                chart_height = Inches(box_height - 1.4)

                pic = slide.shapes.add_picture(chart_buffer, chart_left, chart_top,
                                              width=chart_width, height=chart_height)
                return True

        except Exception as e:
            if self.verbose:
                print(f"⚠️  Grafik oluşturulamadı: {e}")

        return False

    def _add_simple_bar_chart(self, slide, data: Dict, left: float, top: float,
                              width: float, height: float):
        """Add a simple bar chart to the slide"""
        # This would require chart data extraction
        # For now, we'll keep it simple and add placeholder
        pass
